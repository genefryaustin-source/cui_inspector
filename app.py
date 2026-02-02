import sys
import os
import platform
from pathlib import Path
from datetime import datetime
import hashlib
import io
import json
import re
from typing import Dict, Any, List, Tuple

import streamlit as st
import pandas as pd

# ----------------------------
# Ensure project root is importable (Streamlit Cloud safe)
# ----------------------------
ROOT = Path(__file__).resolve().parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# ----------------------------
# Optional OCR dependencies (Cloud-safe)
# ----------------------------
OCR_AVAILABLE = True

try:
    import pytesseract  # type: ignore
except Exception:
    pytesseract = None
    OCR_AVAILABLE = False

try:
    from PIL import Image  # type: ignore
except Exception:
    Image = None

try:
    import pdf2image  # type: ignore
except Exception:
    pdf2image = None
    OCR_AVAILABLE = False

# ----------------------------
# Optional config
# ----------------------------
try:
    from config import TESSERACT_CMD, POPPLER_PATH, OCR_DPI, OCR_LANGUAGE
except Exception:
    TESSERACT_CMD = ""
    POPPLER_PATH = None
    OCR_DPI = 300
    OCR_LANGUAGE = "eng"

if pytesseract and TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

# ----------------------------
# App config
# ----------------------------
st.set_page_config(page_title="CUI Document Inspector", layout="wide")
st.title("📄 CUI Document Inspector")

# =============================================================================
# Clean, indentation-proof sidebar
# =============================================================================
st.sidebar.header("Navigation")
page = st.sidebar.radio("Go to", ["Document Inspector", "System Info"])

st.sidebar.divider()
st.sidebar.markdown("### ✅ Runtime Checks")
st.sidebar.caption(f"OCR available: **{OCR_AVAILABLE}**")

st.sidebar.divider()
st.sidebar.markdown("### ℹ️ System Information")
st.sidebar.write(f"Python: {sys.version.split()[0]}")
st.sidebar.write(f"Platform: {platform.system()} {platform.release()}")

# =============================================================================
# Utilities
# =============================================================================
def now_iso() -> str:
    return datetime.utcnow().isoformat(timespec="seconds") + "Z"


def sha256_bytes(data: bytes) -> str:
    h = hashlib.sha256()
    h.update(data)
    return h.hexdigest()


def clamp(n: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, n))


# =============================================================================
# Extraction
# =============================================================================
def extract_text_from_pdf(uploaded_file) -> str:
    from PyPDF2 import PdfReader

    reader = PdfReader(uploaded_file)
    text_parts: List[str] = []

    for pg in reader.pages:
        t = pg.extract_text() or ""
        if t:
            text_parts.append(t)

    text = "\n".join(text_parts)

    # OCR fallback if scanned/unextractable
    if not text.strip():
        if not OCR_AVAILABLE:
            st.warning("⚠️ PDF appears scanned, but OCR is not available in this deployment.")
            return ""

        try:
            images = pdf2image.convert_from_bytes(
                uploaded_file.getvalue(),
                dpi=OCR_DPI,
                poppler_path=POPPLER_PATH,
            )
            ocr_text = []
            for img in images:
                ocr_text.append(pytesseract.image_to_string(img, lang=OCR_LANGUAGE))
            text = "\n".join(ocr_text)
        except Exception as e:
            st.error(f"OCR failed: {e}")
            return ""

    return text


def extract_text_from_file(uploaded_file) -> str:
    name = (uploaded_file.name or "").lower()

    if name.endswith(".pdf"):
        return extract_text_from_pdf(uploaded_file)

    if name.endswith(".txt"):
        return uploaded_file.getvalue().decode("utf-8", errors="ignore")

    if name.endswith(".docx"):
        from docx import Document
        doc = Document(uploaded_file)
        return "\n".join(p.text for p in doc.paragraphs)

    if name.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        out.append(t)
        return "\n".join(out)

    return ""


# =============================================================================
# Analysis Engine (Step 2)
# Email = sensitive, not CUI (per your instruction)
# =============================================================================
def get_ruleset(name: str) -> Dict[str, Any]:
    """
    Ruleset defines:
      - cui_patterns: patterns considered CUI indicators
      - sensitive_patterns: patterns considered sensitive but not CUI
      - keyword_triggers: words/phrases that increase risk and may imply CUI context
      - weights: scoring weights
    """
    base = {
        "cui_patterns": {
            # Strong CUI indicators (examples)
            "SSN": r"\b\d{3}-\d{2}-\d{4}\b",
            "DOD_ID_10": r"\b\d{10}\b",
            "ITAR_EAR_TERMS": r"\b(ITAR|EAR|USML|DFARS|CUI)\b",
        },
        "sensitive_patterns": {
            # Sensitive but NOT CUI by default (per you)
            "EMAIL": r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}",
            "PHONE_US": r"\b(?:\+1[-.\s]?)?(?:\(\d{3}\)|\d{3})[-.\s]?\d{3}[-.\s]?\d{4}\b",
        },
        "keyword_triggers": [
            "controlled unclassified information",
            "cui",
            "export controlled",
            "itar",
            "ear",
            "dfars",
            "cage code",
            "foia exempt",
            "sensitive but unclassified",
            "for official use only",
            "fouo",
            "proprietary",
            "competition sensitive",
            "spi",
            "ssn",
            "dod id",
        ],
        "weights": {
            "cui_match": 20,
            "sensitive_match": 5,
            "keyword_hit": 8,
            "multi_category_bonus": 10,
        },
    }

    # Slightly stricter variant as a preview (still Step 2)
    if name == "Strict":
        base["weights"]["cui_match"] = 25
        base["weights"]["keyword_hit"] = 10
        base["weights"]["multi_category_bonus"] = 15
    return base


def analyze_text(text: str, ruleset_name: str) -> Dict[str, Any]:
    rules = get_ruleset(ruleset_name)
    t = text or ""
    tlow = t.lower()

    patterns_found: Dict[str, int] = {}
    sensitive_found: Dict[str, int] = {}

    cui_categories: List[str] = []
    sensitive_categories: List[str] = []

    # Pattern matching
    for label, pat in rules["cui_patterns"].items():
        hits = re.findall(pat, t, flags=re.IGNORECASE)
        if hits:
            patterns_found[label] = len(hits)
            cui_categories.append(label)

    for label, pat in rules["sensitive_patterns"].items():
        hits = re.findall(pat, t, flags=re.IGNORECASE)
        if hits:
            sensitive_found[label] = len(hits)
            sensitive_categories.append(label)

    # Keyword triggers
    keyword_hits = []
    for kw in rules["keyword_triggers"]:
        if kw in tlow:
            keyword_hits.append(kw)

    # Score
    score = 0
    score += sum(patterns_found.values()) * rules["weights"]["cui_match"]
    score += sum(sensitive_found.values()) * rules["weights"]["sensitive_match"]
    score += len(keyword_hits) * rules["weights"]["keyword_hit"]

    # bonus for multiple categories
    distinct = len(set(cui_categories))
    if distinct >= 2:
        score += rules["weights"]["multi_category_bonus"]

    # Normalize to 0-100 for UI
    risk_score = int(clamp(score, 0, 100))

    # Risk level mapping
    if risk_score >= 70 or distinct >= 2:
        risk_level = "HIGH"
    elif risk_score >= 30 or distinct == 1:
        risk_level = "MEDIUM"
    else:
        risk_level = "LOW"

    cui_detected = distinct > 0

    # Recommendations (Step 2: lightweight, but useful)
    recommended_actions = []
    if cui_detected:
        recommended_actions += [
            "Mark document with appropriate CUI banner/footer (if applicable).",
            "Confirm handling requirements (storage, transmission, sharing).",
            "Restrict access to need-to-know and record authorized recipients.",
        ]
    if sensitive_found and not cui_detected:
        recommended_actions += [
            "Treat as sensitive: limit distribution and avoid emailing externally.",
            "Redact/replace personal data where feasible.",
        ]
    if not recommended_actions:
        recommended_actions.append("No obvious CUI indicators detected; continue standard document handling.")

    # Placeholder fields for later phases
    evidence_found = []
    missing_requirements = []

    return {
        "ruleset": ruleset_name,
        "cui_detected": cui_detected,
        "risk_score": risk_score,
        "risk_level": risk_level,
        "cui_categories": sorted(list(set(cui_categories))),
        "sensitive_categories": sorted(list(set(sensitive_categories))),
        "patterns_found": patterns_found,           # CUI patterns only
        "sensitive_patterns_found": sensitive_found, # Sensitive patterns only
        "keyword_triggers_hit": keyword_hits,
        "evidence_found": evidence_found,
        "missing_requirements": missing_requirements,
        "recommended_actions": recommended_actions,
    }


def analysis_to_csv_rows(meta: Dict[str, Any], analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows = []
    # Summary row
    rows.append({
        "type": "summary",
        "filename": meta.get("filename", ""),
        "sha256": meta.get("sha256", ""),
        "uploaded_at": meta.get("uploaded_at", ""),
        "ruleset": analysis.get("ruleset", ""),
        "cui_detected": analysis.get("cui_detected", ""),
        "risk_level": analysis.get("risk_level", ""),
        "risk_score": analysis.get("risk_score", ""),
        "cui_categories": ";".join(analysis.get("cui_categories", [])),
        "sensitive_categories": ";".join(analysis.get("sensitive_categories", [])),
        "keyword_triggers_hit": ";".join(analysis.get("keyword_triggers_hit", [])),
    })

    # Detail rows: patterns
    for k, v in (analysis.get("patterns_found") or {}).items():
        rows.append({
            "type": "cui_pattern",
            "filename": meta.get("filename", ""),
            "sha256": meta.get("sha256", ""),
            "key": k,
            "count": v,
        })

    for k, v in (analysis.get("sensitive_patterns_found") or {}).items():
        rows.append({
            "type": "sensitive_pattern",
            "filename": meta.get("filename", ""),
            "sha256": meta.get("sha256", ""),
            "key": k,
            "count": v,
        })

    return rows


# =============================================================================
# Session State
# =============================================================================
if "last_text" not in st.session_state:
    st.session_state.last_text = ""
if "last_meta" not in st.session_state:
    st.session_state.last_meta = {}
if "last_analysis" not in st.session_state:
    st.session_state.last_analysis = None


# =============================================================================
# Pages
# =============================================================================
if page == "Document Inspector":
    st.header("Document Inspection")

    colA, colB = st.columns([1.1, 0.9], gap="large")

    with colA:
        uploaded = st.file_uploader("Upload a document", type=["pdf", "txt", "docx", "pptx"])

        if uploaded:
            extracted_text = extract_text_from_file(uploaded)
            file_bytes = uploaded.getvalue()
            meta = {
                "filename": uploaded.name,
                "size_bytes": uploaded.size,
                "sha256": sha256_bytes(file_bytes),
                "uploaded_at": now_iso(),
            }

            st.session_state.last_text = extracted_text
            st.session_state.last_meta = meta
            st.session_state.last_analysis = None  # reset until user runs analysis

            st.subheader("File Metadata")
            st.json(meta)

            st.subheader("Extracted Text (preview)")
            st.text_area("Preview", extracted_text[:8000], height=260)

        else:
            st.info("Upload a document to begin.")

    with colB:
        st.subheader("Analysis Controls")

        ruleset_name = st.selectbox("Ruleset", ["Basic", "Strict"], index=0)
        st.caption("Email is treated as **sensitive** (not CUI) in this ruleset.")

        can_analyze = bool((st.session_state.last_text or "").strip())
        if not can_analyze:
            st.warning("No extracted text available yet. Upload a document first.")

        run = st.button("▶ Run Analysis", type="primary", disabled=not can_analyze)

        if run:
            analysis = analyze_text(st.session_state.last_text, ruleset_name)
            st.session_state.last_analysis = analysis

        # Render results if available
        if st.session_state.last_analysis:
            analysis = st.session_state.last_analysis
            meta = st.session_state.last_meta or {}

            st.divider()
            st.subheader("Results")

            # Quick summary cards
            c1, c2, c3 = st.columns(3)
            c1.metric("Risk Level", analysis.get("risk_level", ""))
            c2.metric("Risk Score", str(analysis.get("risk_score", "")))
            c3.metric("CUI Detected", "YES" if analysis.get("cui_detected") else "NO")

            if analysis.get("cui_detected"):
                st.warning("CUI indicators detected. Review categories and apply appropriate handling controls.")
            elif analysis.get("sensitive_categories"):
                st.info("Sensitive indicators detected (not CUI). Consider privacy / least-privilege handling.")
            else:
                st.success("No obvious CUI or sensitive indicators detected.")

            with st.expander("Details"):
                st.write("**CUI categories**:", analysis.get("cui_categories", []))
                st.write("**Sensitive categories**:", analysis.get("sensitive_categories", []))
                st.write("**Keyword triggers hit**:", analysis.get("keyword_triggers_hit", []))
                st.write("**CUI patterns found**:", analysis.get("patterns_found", {}))
                st.write("**Sensitive patterns found**:", analysis.get("sensitive_patterns_found", {}))
                st.write("**Recommended actions**:")
                for a in analysis.get("recommended_actions", []):
                    st.write(f"- {a}")

            # Download outputs
            st.divider()
            st.subheader("Download Artifacts")

            payload = {
                "meta": meta,
                "analysis": analysis,
                "generated_at": now_iso(),
            }
            json_bytes = json.dumps(payload, indent=2).encode("utf-8")
            st.download_button(
                "⬇ Download analysis_report.json",
                data=json_bytes,
                file_name="analysis_report.json",
                mime="application/json",
            )

            rows = analysis_to_csv_rows(meta, analysis)
            df = pd.DataFrame(rows)
            csv_bytes = df.to_csv(index=False).encode("utf-8")
            st.download_button(
                "⬇ Download analysis_summary.csv",
                data=csv_bytes,
                file_name="analysis_summary.csv",
                mime="text/csv",
            )

elif page == "System Info":
    st.header("System Information")
    st.write("Python:", sys.version)
    st.write("Platform:", platform.platform())
    st.write("OCR available:", OCR_AVAILABLE)
    st.write("Poppler path:", POPPLER_PATH)
    st.write("Tesseract cmd:", TESSERACT_CMD)


