import streamlit as st

OCR_AVAILABLE = True

try:
    import pytesseract
except Exception:
    pytesseract = None
    OCR_AVAILABLE = False

try:
    import pdf2image
except Exception:
    pdf2image = None
    OCR_AVAILABLE = False

try:
    from config import TESSERACT_CMD, POPPLER_PATH, OCR_DPI, OCR_LANGUAGE
except Exception:
    TESSERACT_CMD = ""
    POPPLER_PATH = None
    OCR_DPI = 300
    OCR_LANGUAGE = "eng"

if pytesseract and TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD


def extract_text_from_pdf(uploaded_file):
    from PyPDF2 import PdfReader
    reader = PdfReader(uploaded_file)
    text = "\n".join(pg.extract_text() or "" for pg in reader.pages)

    if not text.strip() and OCR_AVAILABLE:
        try:
            images = pdf2image.convert_from_bytes(
                uploaded_file.getvalue(),
                dpi=OCR_DPI,
                poppler_path=POPPLER_PATH
            )
            text = "\n".join(
                pytesseract.image_to_string(img, lang=OCR_LANGUAGE)
                for img in images
            )
        except Exception as e:
            st.error(f"OCR failed: {e}")

    return text


def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()

    if name.endswith(".pdf"):
        return extract_text_from_pdf(uploaded_file)
    if name.endswith(".txt"):
        return uploaded_file.getvalue().decode("utf-8", errors="ignore")
    if name.endswith(".docx"):
        from docx import Document
        doc = Document(uploaded_file)
        return "\n".join(p.text for p in doc.paragraphs)
    if name.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    out.append(shape.text.strip())
        return "\n".join(out)
    return ""

